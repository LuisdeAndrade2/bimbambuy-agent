"""
document_loader.py — Carregamento de documentos de múltiplos formatos.

A classe DocumentLoader detecta a extensão do arquivo e usa a biblioteca
adequada para extrair o texto. Cada documento carregado vira um dicionário
no formato:

    {
        "text": "texto completo do documento...",
        "pages": [(1, "texto da página 1"), (2, "texto da página 2")],  # só PDF
        "metadata": {
            "file_name": "politica_reembolsos_devolucoes.pdf",
            "file_path": "data/documents/politica_reembolsos_devolucoes.pdf",
            "extension": ".pdf",
            "modified_date": "2025-08-15 10:30:00",
        },
    }

Formatos suportados: PDF, Word (.docx), Excel (.xlsx), PowerPoint (.pptx),
Markdown (.md), CSV, JSON, HTML e texto puro (.txt).
"""

import os
import csv
import json
import logging
from datetime import datetime

logger = logging.getLogger(__name__)

# Extensões que o carregador sabe processar
SUPPORTED_EXTENSIONS = {
    ".pdf", ".docx", ".xlsx", ".xlsm", ".pptx",
    ".md", ".markdown", ".csv", ".json", ".html", ".htm", ".txt",
}


class DocumentLoader:
    """Carrega documentos de diversos formatos e extrai texto + metadados."""

    # ------------------------------------------------------------------
    # Método principal: detecta a extensão e delega para o loader certo
    # ------------------------------------------------------------------
    def load_document(self, file_path: str) -> dict:
        """
        Detecta a extensão do arquivo e carrega o conteúdo.

        Args:
            file_path: caminho completo do arquivo.

        Returns:
            Dicionário com 'text', 'metadata' e, para PDFs, 'pages'.

        Raises:
            ValueError: se a extensão não for suportada.
        """
        ext = os.path.splitext(file_path)[1].lower()

        if ext not in SUPPORTED_EXTENSIONS:
            raise ValueError(f"Formato não suportado: {ext} ({file_path})")

        logger.info("Carregando documento: %s", file_path)

        if ext == ".pdf":
            return self.load_pdf(file_path)
        elif ext == ".docx":
            return self._load_docx(file_path)
        elif ext in (".xlsx", ".xlsm"):
            return self._load_xlsx(file_path)
        elif ext == ".pptx":
            return self._load_pptx(file_path)
        elif ext in (".md", ".markdown", ".txt"):
            return self._load_plain_text(file_path)
        elif ext == ".csv":
            return self._load_csv(file_path)
        elif ext == ".json":
            return self._load_json(file_path)
        elif ext in (".html", ".htm"):
            return self._load_html(file_path)

        # Nunca deve chegar aqui, mas garantimos um erro claro
        raise ValueError(f"Nenhum loader implementado para: {ext}")

    # ------------------------------------------------------------------
    # PDF — formato prioritário (usa pdfplumber, com fallback pypdf)
    # ------------------------------------------------------------------
    def load_pdf(self, file_path: str) -> dict:
        """
        Extrai texto de um PDF página por página.

        Mantemos o número de cada página para que as citações possam
        indicar "página X" na resposta do agente.
        """
        pages = []
        try:
            # pdfplumber extrai texto com boa fidelidade de layout
            import pdfplumber

            with pdfplumber.open(file_path) as pdf:
                for i, page in enumerate(pdf.pages, start=1):
                    text = page.extract_text() or ""
                    pages.append((i, text))
        except Exception as exc:
            # Fallback: pypdf é mais simples, mas resolve PDFs problemáticos
            logger.warning("pdfplumber falhou (%s). Tentando pypdf...", exc)
            from pypdf import PdfReader

            reader = PdfReader(file_path)
            for i, page in enumerate(reader.pages, start=1):
                text = page.extract_text() or ""
                pages.append((i, text))

        full_text = "\n\n".join(text for _, text in pages)
        return {
            "text": full_text,
            "pages": pages,  # lista de (número_da_página, texto)
            "metadata": self._extract_file_metadata(file_path),
        }

    # ------------------------------------------------------------------
    # Word (.docx)
    # ------------------------------------------------------------------
    def _load_docx(self, file_path: str) -> dict:
        """Extrai parágrafos e tabelas de um arquivo Word."""
        import docx  # python-docx

        doc = docx.Document(file_path)
        parts = [p.text for p in doc.paragraphs if p.text.strip()]

        # Também extrai o conteúdo das tabelas (linha por linha)
        for table in doc.tables:
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                parts.append(" | ".join(cells))

        return {
            "text": "\n".join(parts),
            "metadata": self._extract_file_metadata(file_path),
        }

    # ------------------------------------------------------------------
    # Excel (.xlsx) — cada linha vira uma linha de texto
    # ------------------------------------------------------------------
    def _load_xlsx(self, file_path: str) -> dict:
        """Converte cada aba da planilha em texto (cabeçalho + linhas)."""
        import openpyxl

        wb = openpyxl.load_workbook(file_path, data_only=True, read_only=True)
        parts = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            parts.append(f"=== Planilha: {sheet_name} ===")
            for row in ws.iter_rows(values_only=True):
                values = [str(c) for c in row if c is not None]
                if values:
                    parts.append(" | ".join(values))
        wb.close()

        return {
            "text": "\n".join(parts),
            "metadata": self._extract_file_metadata(file_path),
        }

    # ------------------------------------------------------------------
    # PowerPoint (.pptx) — texto de cada slide
    # ------------------------------------------------------------------
    def _load_pptx(self, file_path: str) -> dict:
        """Extrai o texto de todos os shapes de cada slide."""
        from pptx import Presentation  # python-pptx

        prs = Presentation(file_path)
        parts = []
        for i, slide in enumerate(prs.slides, start=1):
            parts.append(f"=== Slide {i} ===")
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text_frame.text.strip()
                    if text:
                        parts.append(text)

        return {
            "text": "\n".join(parts),
            "metadata": self._extract_file_metadata(file_path),
        }

    # ------------------------------------------------------------------
    # Texto puro / Markdown
    # ------------------------------------------------------------------
    def _load_plain_text(self, file_path: str) -> dict:
        """Lê arquivos .md e .txt diretamente."""
        with open(file_path, "r", encoding="utf-8", errors="replace") as f:
            text = f.read()
        return {
            "text": text,
            "metadata": self._extract_file_metadata(file_path),
        }

    # ------------------------------------------------------------------
    # CSV — cabeçalho + linhas em formato legível
    # ------------------------------------------------------------------
    def _load_csv(self, file_path: str) -> dict:
        """Converte CSV em texto: cada linha vira 'coluna: valor'."""
        parts = []
        with open(file_path, "r", encoding="utf-8", errors="replace", newline="") as f:
            reader = csv.DictReader(f)
            for row in reader:
                line = "; ".join(f"{k}: {v}" for k, v in row.items() if v)
                if line:
                    parts.append(line)

        return {
            "text": "\n".join(parts),
            "metadata": self._extract_file_metadata(file_path),
        }

    # ------------------------------------------------------------------
    # JSON — converte para texto identado (ou percorre listas de objetos)
    # ------------------------------------------------------------------
    def _load_json(self, file_path: str) -> dict:
        """Serializa JSON em texto legível para indexação."""
        with open(file_path, "r", encoding="utf-8", errors="replace") as f:
            data = json.load(f)

        # Se for uma lista de objetos, cada item vira um bloco de texto
        if isinstance(data, list):
            parts = []
            for item in data:
                if isinstance(item, dict):
                    parts.append("; ".join(f"{k}: {v}" for k, v in item.items()))
                else:
                    parts.append(str(item))
            text = "\n".join(parts)
        else:
            # Objeto único: serializa com identação
            text = json.dumps(data, ensure_ascii=False, indent=2)

        return {
            "text": text,
            "metadata": self._extract_file_metadata(file_path),
        }

    # ------------------------------------------------------------------
    # HTML — remove tags e scripts, mantém só o texto
    # ------------------------------------------------------------------
    def _load_html(self, file_path: str) -> dict:
        """Extrai texto visível de um arquivo HTML com BeautifulSoup."""
        from bs4 import BeautifulSoup

        with open(file_path, "r", encoding="utf-8", errors="replace") as f:
            soup = BeautifulSoup(f.read(), "lxml")

        # Remove elementos que não carregam conteúdo útil
        for tag in soup(["script", "style", "noscript"]):
            tag.decompose()

        return {
            "text": soup.get_text(separator="\n"),
            "metadata": self._extract_file_metadata(file_path),
        }

    # ------------------------------------------------------------------
    # Metadados comuns a todos os formatos
    # ------------------------------------------------------------------
    def _extract_file_metadata(self, file_path: str) -> dict:
        """
        Extrai metadados básicos do arquivo: nome, caminho, extensão
        e data da última modificação.
        """
        stat = os.stat(file_path)
        return {
            "file_name": os.path.basename(file_path),
            "file_path": os.path.abspath(file_path),
            "extension": os.path.splitext(file_path)[1].lower(),
            "modified_date": datetime.fromtimestamp(stat.st_mtime).strftime(
                "%Y-%m-%d %H:%M:%S"
            ),
        }
